from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Crear la presentación
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6]

# Colores
blue_color = RGBColor(40, 80, 150)
gray_color = RGBColor(100, 100, 100)

# Slide 1: Portada
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Introducción a GIT"
subtitle.text = "Metodología de Sistemas I\nTecnicatura Universitaria en Programación\nOscar Londero\n24-04-2025"

# Función para agregar slides con título y contenido
def add_content_slide(title_text, content_lines):
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = title_text
    content.text = '\n'.join(content_lines)

# Slide 2: Objetivos
add_content_slide("Objetivos de la Clase", [
    "• Comprender qué es GIT y por qué es importante.",
    "• Instalar GIT en distintos sistemas operativos.",
    "• Conocer y practicar los comandos básicos.",
    "• Aprender el flujo de trabajo estándar con GIT.",
    "• Introducción a GitHub (repositorios remotos)."
])

# Slide 3: ¿Qué es GIT?
add_content_slide("¿Qué es GIT?", [
    "• Sistema de control de versiones distribuido.",
    "• Registra cambios en archivos y coordina el trabajo en equipo.",
    "• Creado por Linus Torvalds en 2005.",
    "Ventajas: historial completo, trabajo colaborativo, integración con GitHub."
])

# Slide 4: Instalación de GIT
add_content_slide("Instalación de GIT", [
    "• Windows: Descargar desde git-scm.com e instalar.",
    "• Linux: sudo apt install git",
    "• macOS: brew install git",
    "• Verificar: git --version"
])

# Slide 5: Configuración Inicial
add_content_slide("Configuración Inicial", [
    'git config --global user.name "Tu Nombre"',
    'git config --global user.email "tuemail@example.com"',
    "Ver configuración: git config --list"
])

# Slide 6: Crear un Repositorio
add_content_slide("Crear un Repositorio", [
    "mkdir mi-proyecto",
    "cd mi-proyecto",
    "git init"
])

# Slide 7: Flujo de Trabajo Básico
add_content_slide("Flujo de Trabajo Básico", [
    "1. git status",
    "2. git add archivo / git add .",
    '3. git commit -m "mensaje"',
    "4. git log"
])

# Slide 8: Comandos Útiles
add_content_slide("Comandos Útiles", [
    "• git diff",
    "• git rm archivo",
    "• git mv viejo nuevo",
    "• git checkout hash",
    "• git reset --hard hash"
])

# Slide 9: GitHub y Repositorio Remoto
add_content_slide("Conectar con GitHub", [
    "• Crear cuenta y repositorio en GitHub.",
    "• git remote add origin URL",
    "• git branch -M main",
    "• git push -u origin main"
])

# Slide 10: Clonar Repositorio
add_content_slide("Clonar un Repositorio", [
    "git clone https://github.com/usuario/repositorio.git"
])

# Slide 11: Buenas Prácticas
add_content_slide("Buenas Prácticas", [
    "• Commits claros y descriptivos.",
    "• Usar .gitignore.",
    "• Hacer pull antes de push.",
    "• Documentar el proyecto (README.md)."
])

# Slide 12: Conclusión
add_content_slide("Conclusión", [
    "• GIT es esencial para el desarrollo moderno.",
    "• Favorece el trabajo en equipo y el orden.",
    "• Dominar GIT es clave como programador."
])

# Slide 13: Práctica Recomendada
add_content_slide("Práctica Recomendada", [
    "• Crear un proyecto personal y versionarlo con GIT.",
    "• Subirlo a GitHub.",
    "• Simular cambios y practicar el flujo completo."
])

# Guardar presentación
pptx_path = "/mnt/data/Presentación_GIT_Metodología_Sistemas_I.pptx"
prs.save(pptx_path)

pptx_path